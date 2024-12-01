import openpyxl
from pptx import Presentation
from pptx.util import Inches

def xlsx_to_pptx(input_excel, output_pptx):
    # Load the Excel workbook
    wb = openpyxl.load_workbook(input_excel)
    sheet = wb.active  # Use the active sheet
    
    # Create a PowerPoint presentation
    prs = Presentation()
    
    # Iterate over the rows in the Excel sheet
    for row_idx, row in enumerate(sheet.iter_rows(values_only=True), start=1):
        if row_idx == 1:
            # Use the first row as the slide title
            title = ", ".join(str(cell) for cell in row if cell is not None)
            continue
        
        # Create a slide
        slide = prs.slides.add_slide(prs.slide_layouts[1])  # Title and Content layout
        title_placeholder = slide.shapes.title
        content_placeholder = slide.placeholders[1]
        
        # Set the slide title
        title_placeholder.text = f"Row {row_idx}"
        
        # Populate the content with the row data
        content = "\n".join(f"Column {i+1}: {cell}" for i, cell in enumerate(row) if cell is not None)
        content_placeholder.text = content
    
    # Save the PowerPoint presentation
    prs.save(output_pptx)
    print(f"PowerPoint presentation saved to {output_pptx}")

# Specify your input Excel file and output PowerPoint file
input_excel = r"/Users/manavk/Documents/secureblox/ml/test-files/Knowledge Base - Questionnaire Sample 2.xlsx"
output_pptx = "k2.pptx"

# Call the function
xlsx_to_pptx(input_excel, output_pptx)
